"""
pptx_importer.py
=================
Reverse of pptx_builder.py: reads an existing, user-supplied .pptx file back
into the editable slide-dict shape the workspace already understands (same
shape /video/render and /video/render/from-pdf return), so a user can upload
a deck they already have and keep editing/re-rendering it here instead of
starting from a blank AI-generated one.

Real-world .pptx files use semantic title/body placeholders, which
`slide.shapes.title` finds directly. But this app's own pptx_builder.py (and
many hand-built decks) use freeform, absolutely-positioned textboxes instead
of placeholders — there is no semantic "title" to query. For those, the
title is identified as the largest-font text shape in the upper portion of
the slide, and tiny shapes pinned to the bottom (footers, page numbers,
brand marks) are excluded from the body content entirely.
"""
from __future__ import annotations

import io
import logging
import re
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)

# Shapes whose top edge falls in the bottom 8% of the slide are treated as
# footer/page-number/brand-mark noise, never real slide content.
_FOOTER_ZONE_FRACTION = 0.92

# Short strings that are near-universally decorative (brand marks, section
# eyebrow labels, confidentiality/date stamps) rather than real slide
# content, regardless of where they sit on the slide.
_NOISE_PATTERNS = (
    re.compile(r"^[A-Z]{1,4}$"),                       # bare brand initials, e.g. "EY"
    re.compile(r"^confidential\b", re.IGNORECASE),
    re.compile(r"^\d{1,3}\s*/\s*\d{1,3}$"),             # "02 / 12" page counters
)


def _is_noise(text: str) -> bool:
    t = text.strip()
    return bool(t) and any(p.match(t) for p in _NOISE_PATTERNS)


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()


def _max_font_size(shape) -> float:
    size = 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                size = max(size, run.font.size.pt)
    return size


def _guess_layout(title_text: str, body_texts: list[str], slide_idx: int, total: int) -> str:
    """Best-effort mapping onto the editor's small layout vocabulary
    (title/content/two_column/quote/metric/closing)."""
    joined = " ".join(body_texts).strip()
    lower = (title_text + " " + joined).lower()
    if slide_idx == 0 and len(joined) < 40:
        return "title"
    if slide_idx == total - 1 and any(w in lower for w in ("thank you", "questions", "conclusion")):
        return "closing"
    if joined.strip().startswith(("“", '"', "'")) or "quote" in title_text.lower():
        return "quote"
    if len(body_texts) >= 2 and all(len(t) < 200 for t in body_texts):
        return "two_column"
    if not joined and title_text:
        return "title"
    return "content"


def parse_pptx_to_slides(file_bytes: bytes) -> list[dict[str, Any]]:
    """Parse an uploaded .pptx into the same slide-dict shape used elsewhere
    in the app: title, subtitle, content, speaker_notes, layout, duration."""
    prs = Presentation(io.BytesIO(file_bytes))
    slide_height = prs.slide_height or 6858000
    total = len(prs.slides)
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides):
        candidates: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = _shape_text(shape)
            if not text or _is_noise(text):
                continue
            top = shape.top if shape.top is not None else 0
            if top >= slide_height * _FOOTER_ZONE_FRACTION:
                continue  # footer / page-number / brand-mark strip — never content
            candidates.append({
                "text": text, "top": top, "font_size": _max_font_size(shape),
                "is_title_placeholder": shape == slide.shapes.title,
            })

        title_text = ""
        body_texts: list[str] = []
        if candidates:
            placeholder_title = next((c for c in candidates if c["is_title_placeholder"] and c["text"]), None)
            title_candidate = placeholder_title or max(candidates, key=lambda c: c["font_size"])
            title_text = title_candidate["text"].splitlines()[0][:120]
            remainder_lines = title_candidate["text"].splitlines()[1:]
            body_texts = (["\n".join(remainder_lines)] if remainder_lines else [])
            body_texts += [c["text"] for c in candidates if c is not title_candidate]
            # Keep original top-to-bottom reading order for the body.
            ordered = sorted(
                (c for c in candidates if c is not title_candidate),
                key=lambda c: c["top"],
            )
            body_texts = [c["text"] for c in ordered]
            if remainder_lines:
                body_texts = ["\n".join(remainder_lines)] + body_texts

        content = "\n".join(
            line if line.lstrip().startswith(("•", "-", "*")) else f"• {line}"
            for text in body_texts
            for line in text.splitlines()
            if line.strip()
        )

        speaker_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            speaker_notes = slide.notes_slide.notes_text_frame.text.strip()

        layout = _guess_layout(title_text, body_texts, idx, total)

        slides.append({
            "title": title_text or f"Slide {idx + 1}",
            "subtitle": "",
            "content": content,
            "speaker_notes": speaker_notes,
            "layout": layout,
            "duration": 25,
        })

    if not slides:
        raise ValueError("No slides found in the uploaded .pptx file")

    return slides
