"""
pptx_builder.py
================
Builds a client-ready, editable .pptx presentation from SDLC artifacts.

Design benchmark: the attached EY corporate template — authentic EY palette
(charcoal 2E2E38 + yellow FFE600 + white), EYInterstate typography, generous
whitespace, card-based information hierarchy, and consulting-grade layouts.

Public entry points (signatures unchanged — nothing that imports this breaks):
  build_pptx(pptx_spec, project_name)                       — from AI agent spec
  build_pptx_from_deck(slides, project_name, theme_id)      — from slide_deck_builder

Supported layouts:
  title · section · agenda · items · stats_grid · kpi_cards · table · chart ·
  two_col · comparison · tech_grid · process · architecture · roadmap ·
  timeline · quote · closing · content (fallback)

All colours / fonts / spacing come from theme_engine.py so the downloadable
deck and the rendered video stay visually identical.
"""
from __future__ import annotations

import io
from .logging_config import get_logger
from datetime import datetime
from pathlib import Path
from typing import Any

from fastapi_agents import theme_engine as TE

logger = get_logger(__name__)

# Kept for backwards-compat: other modules import THEMES from here.
THEMES = TE.THEMES

# Icon names with a hand-drawn vector definition in _icon_glyph (native
# PowerPoint shapes — renders identically on every machine, no font/emoji
# coverage issues). Names outside this set but in TE.ICONS still resolve via
# the same vector path when they alias one of these (see _icon_glyph's
# generic fallback); anything else falls back to a short text label.
_VECTOR_ICON_NAMES = {
    "shield", "lock", "gear", "cloud", "database", "chart", "growth", "rocket",
    "users", "user", "api", "flow", "network", "check", "test", "warning",
    "target", "clock", "money", "roi", "idea", "brain", "layers", "doc",
    "code", "flag", "eye",
}


def _rgb(hex_color: str) -> tuple[int, int, int]:
    h = str(hex_color).lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _emu(inches: float) -> int:
    return int(inches * 914400)


# ── Main entry points ─────────────────────────────────────────────────────────

def build_pptx(pptx_spec: dict[str, Any], project_name: str = "SDLC Project") -> bytes:
    """Build PPTX from an AI-agent pptx_spec. Rich slides in the spec are mapped
    onto the deck layout vocabulary; if the spec is thin we fall back to the
    full artifact-driven deck builder so the output is never empty."""
    from fastapi_agents.slide_deck_builder import build_deck

    theme_id = TE.DEFAULT_THEME
    theme_field = pptx_spec.get("theme") if isinstance(pptx_spec, dict) else None
    if isinstance(theme_field, dict):
        theme_id = theme_field.get("theme_id") or theme_field.get("id") or theme_id

    slides_data = pptx_spec.get("slides", []) if isinstance(pptx_spec, dict) else []
    deck_slides: list[dict] = []
    for i, s in enumerate(slides_data):
        if not isinstance(s, dict):
            continue
        # If the spec already speaks deck vocabulary, use it directly.
        if s.get("layout"):
            deck_slides.append(s)
            continue
        mapped = _map_ai_slide_to_deck(s, i, project_name)
        if mapped:
            deck_slides.append(mapped)

    if not deck_slides:
        deck_slides = build_deck([], project_name)

    return _build_from_slides(deck_slides, project_name, theme_id)


def build_pptx_from_deck(slides: list[dict], project_name: str = "SDLC Project",
                         theme_id: str = "ey") -> bytes:
    """Build PPTX directly from slide_deck_builder slides."""
    return _build_from_slides(slides, project_name, theme_id)


# Map the Storytelling/Review agent's layout hints onto the deck vocabulary.
_LAYOUT_ALIASES = {
    "title_slide": "title", "title": "title", "section": "section",
    "divider": "section", "agenda": "agenda", "summary": "kpi_cards",
    "data": "kpi_cards", "kpi": "kpi_cards", "kpi_cards": "kpi_cards",
    "stats": "stats_grid", "stats_grid": "stats_grid", "metrics": "kpi_cards",
    "table": "table", "chart": "chart", "bar": "chart", "progress": "chart",
    "two_col": "two_col", "two_content": "two_col", "two_column": "two_col",
    "comparison": "comparison", "before_after": "comparison",
    "tech": "tech_grid", "tech_grid": "tech_grid", "technology": "tech_grid",
    "process": "process", "workflow": "process", "steps": "process",
    "architecture": "architecture", "roadmap": "roadmap", "timeline": "timeline",
    "quote": "quote", "closing": "closing", "conclusion": "closing",
    "items": "items", "content": "content", "bullet-list": "items",
    "hero": "hero", "hero_image": "hero", "vision": "hero", "statement": "hero",
}


def _hint_to_layout(s: dict) -> str:
    for key in (s.get("visual_suggestions"), s.get("slide_type"),
                (s.get("pptx_layout", {}) or {}).get("layout_type")):
        if not key:
            continue
        k = str(key).strip().lower().replace(" ", "_")
        if k in _LAYOUT_ALIASES:
            return _LAYOUT_ALIASES[k]
        for token, layout in _LAYOUT_ALIASES.items():
            if token in k:
                return layout
    return ""


def _map_ai_slide_to_deck(s: dict, idx: int, project_name: str) -> dict | None:
    """Convert one Storytelling/Review agent slide into a rich deck slide,
    honouring the planned layout and best-effort structuring the AI's bullets
    into the target layout's shape."""
    content = s.get("content", {}) or {}
    if not isinstance(content, dict):
        content = {"body_text": str(content)}
    bullets = content.get("bullets", []) or []
    body = content.get("body_text", "") or ""
    data_pts = content.get("data_points", []) or []
    title = s.get("title", "") or (project_name if idx == 0 else "")
    subtitle = s.get("subtitle", "")
    notes = s.get("speaker_notes", "")
    layout = _hint_to_layout(s)
    if idx == 0 and layout in ("", "content", "items"):
        layout = "title"

    # Carried through unchanged for every layout branch below (all spread
    # **base) — image_prompt is the LLM-authored per-slide image prompt
    # (empty by default, doesn't affect anything if unset); pptx_layout
    # carries the LLM's optional typography/color overrides, applied by
    # _header() on top of the theme defaults when present.
    base = {
        "title": title, "subtitle": subtitle, "speaker_notes": notes,
        "image_prompt": s.get("image_prompt", "") or "",
        "pptx_layout": s.get("pptx_layout") or {},
    }

    def _split(b):
        parts = str(b).split(" — ") if " — " in str(b) else str(b).split(": ", 1)
        return (parts[0].strip(), parts[1].strip() if len(parts) > 1 else "")

    if layout == "title":
        return {**base, "layout": "title", "badge": "EXECUTIVE BRIEF",
                "subtitle": subtitle or "AI-Autonomous Software Delivery"}
    if layout == "section":
        return {**base, "layout": "section", "number": f"{idx:02d}"}
    if layout == "closing":
        items = [{"icon": "check", "title": _split(b)[0], "body": _split(b)[1]} for b in bullets[:4]]
        return {**base, "layout": "closing", "title": title or "Thank You", "items": items}
    if layout == "quote":
        return {**base, "layout": "quote", "quote": body or title, "attribution": subtitle}
    if layout == "hero":
        content_text = body or (" — ".join(bullets[:2]) if bullets else "")
        return {**base, "layout": "hero", "content": content_text}
    if layout in ("kpi_cards", "stats_grid") or (data_pts and not layout):
        stats = [{"value": str(d.get("value", "")), "label": d.get("label", ""),
                  "sub": d.get("context", "")} for d in data_pts[:4]]
        if not stats and bullets:
            stats = [{"value": f"{i+1:02d}", "label": _split(b)[0], "sub": _split(b)[1]}
                     for i, b in enumerate(bullets[:4])]
        return {**base, "layout": layout or "kpi_cards", "stats": stats}
    if layout == "agenda":
        items = [{"title": _split(b)[0], "body": _split(b)[1]} for b in bullets[:6]]
        return {**base, "layout": "agenda", "items": items}
    if layout == "process":
        steps = [{"title": _split(b)[0], "body": _split(b)[1]} for b in bullets[:5]]
        return {**base, "layout": "process", "steps": steps}
    if layout == "architecture":
        layers = [{"name": _split(b)[0], "desc": _split(b)[1], "icon": "layers"} for b in bullets[:5]]
        return {**base, "layout": "architecture", "layers": layers}
    if layout == "roadmap" or layout == "timeline":
        evs = [{"title": _split(b)[0], "desc": _split(b)[1], "date": f"Phase {i+1}"}
               for i, b in enumerate(bullets[:5])]
        return {**base, "layout": layout, "timeline": evs}
    if layout == "tech_grid":
        techs = [{"layer": _split(b)[0], "tech": _split(b)[1] or _split(b)[0], "icon": "gear"}
                 for b in bullets[:8]]
        return {**base, "layout": "tech_grid", "tech_items": techs}
    if layout == "comparison":
        half = (len(bullets) + 1) // 2
        return {**base, "layout": "comparison",
                "left": {"header": "Today", "points": bullets[:half]},
                "right": {"header": "With the Platform", "points": bullets[half:]}}
    if layout == "two_col":
        half = (len(bullets) + 1) // 2
        return {**base, "layout": "two_col",
                "left_header": "Challenge", "left_content": "\n".join(bullets[:half]),
                "right_header": "Approach", "right_content": "\n".join(bullets[half:])}
    if layout == "chart":
        data = [{"label": _split(b)[0], "value": 70 + (i * 5) % 30, "max": 100}
                for i, b in enumerate(bullets[:6])]
        return {**base, "layout": "chart", "chart": {"data": data}}
    if layout == "table" and bullets:
        rows = [[_split(b)[0], _split(b)[1]] for b in bullets[:10]]
        return {**base, "layout": "table", "table": {"headers": ["Item", "Detail"], "rows": rows}}
    # Default: card list if we have bullets, else a content panel
    if bullets:
        return {**base, "layout": "items",
                "items": [{"icon": "check", "title": _split(b)[0], "body": _split(b)[1]}
                          for b in bullets[:6]]}
    if body:
        return {**base, "layout": "content", "content": body}
    return {**base, "layout": "content", "content": subtitle or title}


# ── Core builder ──────────────────────────────────────────────────────────────

def _build_from_slides(slides: list[dict], project_name: str, theme_id: str = "ey") -> bytes:
    # Local-first design engine pass: every slide that arrives with a generic
    # "items"/"content" layout gets upgraded to the visual layout its content
    # actually calls for (KPI cards, process flow, comparison, roadmap,
    # architecture) and its bullets trimmed to visual-first phrases, BEFORE
    # rendering. pptx_builder itself stays purely a rendering/export layer.
    # Never blocks: falls back to the original slides on any failure.
    try:
        from .design_engine import design_slides
        slides = design_slides(slides)
    except Exception as exc:
        logger.debug("[PptxBuilder] design engine pass skipped: %s", exc)

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.oxml.ns import qn
    except ImportError:
        raise RuntimeError("python-pptx not installed: pip install python-pptx")

    theme = TE.get_theme(theme_id)
    L = TE.LAYOUT
    F_HEAD = TE.FONT_STACK["pptx_heading"]
    F_BODY = TE.FONT_STACK["pptx_body"]

    C = {k: RGBColor(*_rgb(v)) for k, v in theme.items()
         if isinstance(v, str) and len(str(v)) == 6 and all(ch in "0123456789ABCDEFabcdef" for ch in str(v))}

    def col(name: str) -> RGBColor:
        return C.get(name, RGBColor(*_rgb(theme.get(name, "2E2E38"))))

    is_dark = theme.get("mode") == "dark"

    def _luminance(hexname: str) -> float:
        r, g, b = _rgb(theme.get(hexname, "FFFFFF"))
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    # Hero/dark-background slides (title/section/quote/closing) use theme
    # "accent" for badges, kickers, and subtitle text. That's correct for
    # themes where accent is a bright color (EY yellow, government gold,
    # startup violet, healthcare teal) — but McKinsey's "accent" is a near-
    # black navy (051C2C), designed to read on a *light* background, and is
    # nearly invisible against its own dark hero gradient. Fall back to
    # "info" (a bright blue in every theme) whenever accent's luminance is
    # too low to read on a dark surface.
    _HERO_ACCENT = "accent" if _luminance("accent") > 60 else "info"

    def hero_accent() -> RGBColor:
        return col(_HERO_ACCENT)

    prs = Presentation()
    prs.slide_width = Emu(_emu(TE.CANVAS["w"]))
    prs.slide_height = Emu(_emu(TE.CANVAS["h"]))
    blank = prs.slide_layouts[6]

    # ── primitive helpers ─────────────────────────────────────────────────
    def _bg(slide, hexname="bg"):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = col(hexname)

    def _lighten(hexname: str, amount: float = 0.18) -> RGBColor:
        """Shift a theme color toward white by `amount` (0-1). Used to derive
        a guaranteed-visible second gradient stop programmatically, since
        some themes' semantic "accent2" is identical to (or very close to)
        "bg3" — a fixed named-color pairing would render as a flat, gradient-
        less wash on those themes (confirmed: government and mckinsey both
        have accent2 == bg3)."""
        r, g, b = _rgb(theme.get(hexname, "2E2E38"))
        r = int(r + (255 - r) * amount)
        g = int(g + (255 - g) * amount)
        b = int(b + (255 - b) * amount)
        return RGBColor(r, g, b)

    def _bg_gradient(slide, hex_from="bg3", angle=45, amount=0.22):
        """Full-bleed two-stop gradient background — a subtle diagonal wash
        from the theme's hero color to a programmatically lightened variant
        of itself (guaranteed visible contrast on every theme, unlike a fixed
        named second color) — for hero-moment slides (title/section/quote/
        closing). Implemented as a borderless rectangle behind everything
        else, since python-pptx's slide.background has no gradient API —
        shapes do."""
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                       Emu(_emu(TE.CANVAS["w"])), Emu(_emu(TE.CANVAS["h"])))
        rect.line.fill.background()
        rect.shadow.inherit = False
        rect.fill.gradient()
        stops = rect.fill.gradient_stops
        stops[0].color.rgb = col(hex_from)
        stops[0].position = 0.0
        stops[-1].color.rgb = _lighten(hex_from, amount)
        stops[-1].position = 1.0
        try:
            rect.fill.gradient_angle = angle
        except Exception:
            pass
        _flatten(rect)
        # Send to back so subsequent shapes/text draw on top.
        spTree = slide.shapes._spTree
        spTree.remove(rect._element)
        spTree.insert(2, rect._element)
        return rect

    def _set_alpha(shp, pct_opaque: int) -> None:
        """Shape transparency isn't exposed by python-pptx's high-level API —
        set it via the underlying <a:solidFill><a:srgbClr><a:alpha> XML.
        pct_opaque: 0-100 (100 = fully opaque)."""
        try:
            sp = shp.fill.fore_color._xFill
            srgb = sp.find(qn("a:srgbClr"))
            if srgb is None:
                return
            alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(pct_opaque * 1000))})
            srgb.append(alpha)
        except Exception:
            pass

    def _decorative_blobs(slide, accent_hex="accent", accent2_hex="on_dark", corner="br"):
        """Soft, semi-transparent overlapping circles in one corner — the
        understated geometric-accent treatment Gamma-style decks use to keep
        a hero slide from reading as a flat block of color. Never overlaps
        the title-safe content area (kept to the outer ~25% of the canvas).
        Supports all four corners: "br" | "tr" | "bl" | "tl"."""
        w, h = TE.CANVAS["w"], TE.CANVAS["h"]
        cx = w - 0.4 if corner in ("br", "tr") else 0.4
        cy = h - 0.3 if corner in ("br", "bl") else 0.3
        specs = [(2.6, 18), (1.7, 12), (1.0, 26)]
        for size, alpha in specs:
            o = _shape(slide, MSO_SHAPE.OVAL, cx - size / 2, cy - size / 2, size, size, fill=col(accent_hex))
            _set_alpha(o, alpha)
        # One small solid accent dot near the blob cluster for a crisp focal point.
        dot = _shape(slide, MSO_SHAPE.OVAL, cx - 0.9, cy - 1.6, 0.22, 0.22, fill=col(accent2_hex))
        _set_alpha(dot, 55)

    def _flatten(shp):
        """Force a flat, shadow-free shape (authentic EY/McKinsey decks are
        flat). python-pptx autoshapes carry a <p:style> with an effectRef that
        LibreOffice/PowerPoint resolve to a theme drop-shadow; removing the
        style node and appending an empty <a:effectLst/> guarantees no shadow."""
        el = shp._element
        spPr = el.spPr
        for style in el.findall(qn("p:style")):
            el.remove(style)
        for eff in spPr.findall(qn("a:effectLst")):
            spPr.remove(eff)
        spPr.append(spPr.makeelement(qn("a:effectLst"), {}))

    def _shape(slide, kind, l, t, w, h, fill=None, line_rgb=None, line_w=1.0):
        shp = slide.shapes.add_shape(kind, Emu(_emu(l)), Emu(_emu(t)), Emu(_emu(w)), Emu(_emu(h)))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line_rgb is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line_rgb
            shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        _flatten(shp)
        return shp

    def _rect(slide, l, t, w, h, fill=None, line_rgb=None, line_w=1.0):
        return _shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, fill, line_rgb, line_w)

    def _card(slide, l, t, w, h, fill="bg2", border="hairline"):
        """Rounded card with subtle border — the core information container."""
        shp = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
                     fill=col(fill), line_rgb=col(border) if border else None, line_w=0.75)
        try:  # tame the corner radius (default is huge)
            shp.adjustments[0] = 0.045
        except Exception:
            pass
        return shp

    def _card_glass(slide, l, t, w, h, *, tint="on_dark", alpha=16, border_alpha=35):
        """Semi-transparent 'glass' card — for content that needs to sit on
        top of a hero gradient/image rather than a flat theme background.
        Uses the same alpha-XML technique as _decorative_blobs/_set_alpha."""
        shp = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h,
                     fill=col(tint), line_rgb=col(tint), line_w=0.75)
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
        _set_alpha(shp, alpha)
        try:
            line_srgb = shp.line.color._xFill.find(qn("a:srgbClr"))
            if line_srgb is not None:
                alpha_el = line_srgb.makeelement(qn("a:alpha"), {"val": str(int(border_alpha * 1000))})
                line_srgb.append(alpha_el)
        except Exception:
            pass
        return shp

    def _tb(slide, text, l, t, w, h, size=18, bold=False, color=None, italic=False,
            align=PP_ALIGN.LEFT, font=None, anchor=MSO_ANCHOR.TOP, spacing=1.06,
            wrap=True, shrink=False):
        txb = slide.shapes.add_textbox(Emu(_emu(l)), Emu(_emu(t)), Emu(_emu(w)), Emu(_emu(h)))
        tf = txb.text_frame
        tf.word_wrap = wrap
        try:
            tf.vertical_anchor = anchor
            tf.margin_left = tf.margin_right = Emu(_emu(0.04))
            tf.margin_top = tf.margin_bottom = Emu(_emu(0.02))
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = str(text)
        run.font.name = font or F_HEAD
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color if color is not None else col("body_text")
        return txb

    def _bullets(slide, items, l, t, w, h, size=14, color=None, gap=6,
                 marker="—", marker_color=None, bold_lead=False):
        txb = slide.shapes.add_textbox(Emu(_emu(l)), Emu(_emu(t)), Emu(_emu(w)), Emu(_emu(h)))
        tf = txb.text_frame
        tf.word_wrap = True
        body_c = color if color is not None else col("body_text")
        mk_c = marker_color if marker_color is not None else col("accent2")
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(gap if i else 0)
            p.line_spacing = 1.12
            if marker:
                mr = p.add_run(); mr.text = f"{marker}  "
                mr.font.name = F_HEAD; mr.font.size = Pt(size); mr.font.bold = True
                mr.font.color.rgb = mk_c
            r = p.add_run(); r.text = str(item).lstrip("•-–—▪ ").strip()
            r.font.name = F_BODY; r.font.size = Pt(size); r.font.color.rgb = body_c
            r.font.bold = bold_lead
        return txb

    def _icon_glyph(slide, name, l, t, size_in, ink):
        """Draw a small flat vector icon (native PowerPoint shapes, not a
        font glyph) so it renders identically on every machine/Office version
        — no missing-glyph tofu boxes, no emoji-font inconsistency. Covers the
        icon vocabulary used across layouts; anything outside that set falls
        back to a clean check/dot mark rather than raw Unicode text."""
        cx, cy = l + size_in / 2, t + size_in / 2
        s = size_in
        lw = max(s * 11, 1.75)

        def line(x1, y1, x2, y2, w=lw, color=None):
            ln = slide.shapes.add_connector(1, Emu(_emu(x1)), Emu(_emu(y1)), Emu(_emu(x2)), Emu(_emu(y2)))
            ln.line.color.rgb = color or ink
            ln.line.width = Pt(w)
            ln.shadow.inherit = False
            return ln

        def oval(x, y, w, h, fill=None, outline=None, ow=lw):
            o = _shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill=fill, line_rgb=outline, line_w=ow)
            return o

        def rect(x, y, w, h, fill=None, outline=None, ow=lw, rounded=False):
            kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
            r = _shape(slide, kind, x, y, w, h, fill=fill, line_rgb=outline, line_w=ow)
            return r

        n = (name or "").lower().strip()
        if n == "shield":
            sh = _shape(slide, MSO_SHAPE.PENTAGON, cx - s * 0.28, cy - s * 0.32, s * 0.56, s * 0.62,
                       fill=None, line_rgb=ink, line_w=lw)
            try:
                sh.rotation = 180
            except Exception:
                pass
            line(cx - s * 0.12, cy - s * 0.02, cx - s * 0.02, cy + s * 0.1, w=lw)
            line(cx - s * 0.02, cy + s * 0.1, cx + s * 0.16, cy - s * 0.14, w=lw)
        elif n == "lock":
            rect(cx - s * 0.22, cy - s * 0.04, s * 0.44, s * 0.32, outline=ink, rounded=True)
            arc = _shape(slide, MSO_SHAPE.BLOCK_ARC, cx - s * 0.16, cy - s * 0.32, s * 0.32, s * 0.32,
                        fill=None, line_rgb=ink, line_w=lw)
            oval(cx - s * 0.04, cy + s * 0.06, s * 0.08, s * 0.08, fill=ink)
        elif n == "gear":
            oval(cx - s * 0.22, cy - s * 0.22, s * 0.44, s * 0.44, outline=ink)
            oval(cx - s * 0.08, cy - s * 0.08, s * 0.16, s * 0.16, fill=ink)
            for ang in (0, 45, 90, 135):
                import math
                dx, dy = s * 0.30 * math.cos(math.radians(ang)), s * 0.30 * math.sin(math.radians(ang))
                line(cx - dx, cy - dy, cx + dx, cy + dy, w=lw)
        elif n == "cloud":
            oval(cx - s * 0.26, cy - s * 0.06, s * 0.28, s * 0.24, fill=ink)
            oval(cx - s * 0.08, cy - s * 0.16, s * 0.30, s * 0.30, fill=ink)
            oval(cx + s * 0.10, cy - s * 0.04, s * 0.22, s * 0.20, fill=ink)
            rect(cx - s * 0.26, cy + s * 0.06, s * 0.58, s * 0.14, fill=ink)
        elif n == "database":
            for dy in (-0.20, -0.02, 0.16):
                oval(cx - s * 0.24, cy + s * dy, s * 0.48, s * 0.14, outline=ink, ow=lw * 0.8)
            line(cx - s * 0.24, cy - s * 0.13, cx - s * 0.24, cy + s * 0.23, w=lw * 0.8)
            line(cx + s * 0.24, cy - s * 0.13, cx + s * 0.24, cy + s * 0.23, w=lw * 0.8)
        elif n in ("chart", "growth"):
            base_y = cy + s * 0.22
            for i, h in enumerate((0.18, 0.32, 0.24, 0.4)):
                bx = cx - s * 0.30 + i * s * 0.20
                rect(bx, base_y - s * h, s * 0.14, s * h, fill=ink)
        elif n == "rocket":
            tri = _shape(slide, MSO_SHAPE.ISOCELES_TRIANGLE, cx - s * 0.12, cy - s * 0.30, s * 0.24, s * 0.5,
                        fill=ink)
            oval(cx - s * 0.28, cy + s * 0.14, s * 0.16, s * 0.16, fill=ink)
            oval(cx + s * 0.12, cy + s * 0.14, s * 0.16, s * 0.16, fill=ink)
        elif n in ("users", "user"):
            oval(cx - s * 0.20, cy - s * 0.22, s * 0.22, s * 0.22, fill=ink)
            oval(cx + s * 0.02, cy - s * 0.14, s * 0.18, s * 0.18, fill=ink)
            _shape(slide, MSO_SHAPE.OVAL, cx - s * 0.28, cy + s * 0.02, s * 0.34, s * 0.28, fill=ink)
            _shape(slide, MSO_SHAPE.OVAL, cx + s * 0.02, cy + s * 0.08, s * 0.28, s * 0.24, fill=ink)
        elif n in ("api", "flow", "network"):
            oval(cx - s * 0.30, cy - s * 0.08, s * 0.16, s * 0.16, fill=ink)
            oval(cx + s * 0.14, cy - s * 0.08, s * 0.16, s * 0.16, fill=ink)
            oval(cx - s * 0.08, cy + s * 0.14, s * 0.16, s * 0.16, fill=ink)
            line(cx - s * 0.22, cy, cx + s * 0.22, cy, w=lw * 0.8)
            line(cx - s * 0.22, cy, cx, cy + s * 0.22, w=lw * 0.8)
        elif n in ("check", "test"):
            line(cx - s * 0.18, cy, cx - s * 0.03, cy + s * 0.16, w=lw * 1.3)
            line(cx - s * 0.03, cy + s * 0.16, cx + s * 0.24, cy - s * 0.16, w=lw * 1.3)
        elif n == "warning":
            _shape(slide, MSO_SHAPE.ISOCELES_TRIANGLE, cx - s * 0.28, cy - s * 0.24, s * 0.56, s * 0.5,
                  fill=None, line_rgb=ink, line_w=lw)
            line(cx, cy - s * 0.02, cx, cy + s * 0.08, w=lw * 1.2)
            oval(cx - s * 0.018, cy + s * 0.15, s * 0.036, s * 0.036, fill=ink)
        elif n in ("target",):
            oval(cx - s * 0.26, cy - s * 0.26, s * 0.52, s * 0.52, outline=ink)
            oval(cx - s * 0.13, cy - s * 0.13, s * 0.26, s * 0.26, outline=ink, ow=lw * 0.8)
            oval(cx - s * 0.04, cy - s * 0.04, s * 0.08, s * 0.08, fill=ink)
        elif n == "clock":
            oval(cx - s * 0.26, cy - s * 0.26, s * 0.52, s * 0.52, outline=ink)
            line(cx, cy, cx, cy - s * 0.16, w=lw * 0.9)
            line(cx, cy, cx + s * 0.12, cy + s * 0.04, w=lw * 0.9)
        elif n in ("money", "roi"):
            oval(cx - s * 0.24, cy - s * 0.24, s * 0.48, s * 0.48, outline=ink)
            _tb(slide, "$", cx - s * 0.13, cy - s * 0.20, s * 0.26, s * 0.4, size=int(s * 20),
              bold=True, color=ink, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif n in ("idea", "brain"):
            oval(cx - s * 0.18, cy - s * 0.24, s * 0.36, s * 0.34, outline=ink)
            line(cx - s * 0.08, cy + s * 0.10, cx + s * 0.08, cy + s * 0.10, w=lw)
            line(cx - s * 0.06, cy + s * 0.18, cx + s * 0.06, cy + s * 0.18, w=lw * 0.8)
        elif n in ("layers", "doc"):
            for i, dy in enumerate((-0.14, 0.0, 0.14)):
                _shape(slide, MSO_SHAPE.PARALLELOGRAM, cx - s * 0.26, cy + s * dy - s * 0.06,
                      s * 0.52, s * 0.12, fill=(ink if i == 1 else None),
                      line_rgb=(None if i == 1 else ink))
        elif n in ("code",):
            line(cx - s * 0.06, cy - s * 0.18, cx - s * 0.22, cy, w=lw)
            line(cx - s * 0.22, cy, cx - s * 0.06, cy + s * 0.18, w=lw)
            line(cx + s * 0.06, cy - s * 0.18, cx + s * 0.22, cy, w=lw)
            line(cx + s * 0.22, cy, cx + s * 0.06, cy + s * 0.18, w=lw)
        elif n in ("flag",):
            line(cx - s * 0.2, cy - s * 0.26, cx - s * 0.2, cy + s * 0.26, w=lw)
            _shape(slide, MSO_SHAPE.RECTANGLE, cx - s * 0.2, cy - s * 0.24, s * 0.4, s * 0.2, fill=ink)
        elif n in ("eye",):
            _shape(slide, MSO_SHAPE.OVAL, cx - s * 0.28, cy - s * 0.14, s * 0.56, s * 0.28, fill=None,
                  line_rgb=ink, line_w=lw)
            oval(cx - s * 0.07, cy - s * 0.07, s * 0.14, s * 0.14, fill=ink)
        else:
            # Generic mark: clean checkmark reads better than a missing-glyph box.
            line(cx - s * 0.18, cy, cx - s * 0.03, cy + s * 0.16, w=lw * 1.3)
            line(cx - s * 0.03, cy + s * 0.16, cx + s * 0.24, cy - s * 0.16, w=lw * 1.3)

    def _icon_chip(slide, l, t, size_in, icon_name, fill="accent", glyph_color="on_accent"):
        """`icon_name` is a semantic icon name (e.g. "shield", "check",
        "database") — resolved to a hand-drawn vector icon so it renders
        identically everywhere. A short literal label (e.g. "vs") that isn't
        a known icon name is drawn as text instead."""
        chip = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, size_in, size_in, fill=col(fill))
        try:
            chip.adjustments[0] = 0.28
        except Exception:
            pass
        name = str(icon_name).lower().strip()
        if name in TE.ICONS or name in _VECTOR_ICON_NAMES:
            _icon_glyph(slide, name, l, t, size_in, col(glyph_color))
        else:
            _tb(slide, str(icon_name), l, t + size_in * 0.06, size_in, size_in * 0.82,
                size=int(size_in * 22), bold=True, color=col(glyph_color),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return chip

    def _notes(slide, text):
        if text:
            slide.notes_slide.notes_text_frame.text = str(text)

    def _picture(slide, path, l, t, w, h):
        """Place an image (e.g. an auto-generated diagram) fit within the
        given box, centered, preserving aspect ratio. No-ops silently if the
        file is missing/unreadable so a bad diagram path never breaks the deck."""
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(path) as im:
                iw, ih = im.size
            box_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > box_ratio:
                pw, ph = w, w / img_ratio
            else:
                ph, pw = h, h * img_ratio
            px, py = l + (w - pw) / 2, t + (h - ph) / 2
            slide.shapes.add_picture(str(path), Emu(_emu(px)), Emu(_emu(py)),
                                     width=Emu(_emu(pw)), height=Emu(_emu(ph)))
        except Exception as exc:
            logger.debug("[PptxBuilder] diagram image skipped (%s): %s", path, exc)

    def _picture_cover(slide, path, l, t, w, h):
        """Like _picture() but fills the box completely (cropping overflow)
        instead of letterboxing with empty bars — for hero-image regions
        where empty space reads as unpolished. Uses python-pptx's crop_*
        properties, which trim the *displayed* portion of the image without
        touching the source file. Returns the Picture shape, or None if the
        image couldn't be placed (never raises — a bad/missing hero image
        just means that slide keeps its gradient-only look)."""
        try:
            from PIL import Image as _PILImage
            with _PILImage.open(path) as im:
                iw, ih = im.size
            box_ratio = w / h
            img_ratio = iw / ih
            pic = slide.shapes.add_picture(str(path), Emu(_emu(l)), Emu(_emu(t)),
                                            width=Emu(_emu(w)), height=Emu(_emu(h)))
            if img_ratio > box_ratio:
                visible_frac = box_ratio / img_ratio
                crop = (1 - visible_frac) / 2
                pic.crop_left = crop
                pic.crop_right = crop
            else:
                visible_frac = img_ratio / box_ratio
                crop = (1 - visible_frac) / 2
                pic.crop_top = crop
                pic.crop_bottom = crop
            return pic
        except Exception as exc:
            logger.debug("[PptxBuilder] hero image skipped (%s): %s", path, exc)
            return None

    def _hero_wash(slide, s, *, width_frac=0.3, alpha=14):
        """Low-alpha corner illustration behind an otherwise-light content
        layout (e.g. KPI cards) — 'industry images combined with KPI cards':
        the cards stay fully opaque/legible in front, the image just adds
        atmosphere in the header/margin area behind them. No-ops if no
        hero_image is present."""
        hero_image = s.get("hero_image")
        if not (hero_image and Path(hero_image).exists()):
            return
        w = TE.CANVAS["w"] * width_frac
        pic = _picture_cover(slide, hero_image, TE.CANVAS["w"] - w, 0, w, TE.CANVAS["h"])
        if pic is not None:
            try:
                pic.line.fill.background()
            except Exception:
                pass
            # python-pptx doesn't expose picture alpha directly; a
            # transparent scrim rectangle on top approximates the same
            # "faded into the background" effect cheaply and reliably.
            scrim = _rect(slide, TE.CANVAS["w"] - w, 0, w, TE.CANVAS["h"], fill=col("bg"))
            _set_alpha(scrim, 100 - alpha)

    def _kicker(slide, text, l, t):
        """Small uppercase eyebrow label above a title (consulting hallmark)."""
        if not text:
            return
        _rect(slide, l, t + 0.02, 0.34, 0.045, fill=col("accent"))
        _tb(slide, str(text).upper(), l + 0.42, t - 0.14, 6.5, 0.32,
            size=11, bold=True, color=col("accent3"), spacing=1.0)

    def _footer(slide, n, total, title=""):
        _rect(slide, 0, L["footer_y"], TE.CANVAS["w"], TE.CANVAS["h"] - L["footer_y"],
              fill=col("bg"))
        _rect(slide, L["margin_x"], L["footer_y"], TE.CANVAS["w"] - 2 * L["margin_x"],
              0.012, fill=col("hairline"))
        _tb(slide, "EY  ·  Autonomous SDLC Studio", L["margin_x"], L["footer_y"] + 0.06,
            5.5, 0.3, size=9, color=col("muted"))
        if title:
            _tb(slide, title[:64], 4.5, L["footer_y"] + 0.06, 4.3, 0.3,
                size=9, color=col("muted"), align=PP_ALIGN.CENTER)
        _tb(slide, f"{n:02d} / {total:02d}", TE.CANVAS["w"] - L["margin_x"] - 1.5,
            L["footer_y"] + 0.06, 1.5, 0.3, size=9, color=col("muted"), align=PP_ALIGN.RIGHT)

    def _header(slide, title, subtitle, kicker=""):
        """Standard content-slide header, McKinsey/EY style: a topic KICKER, then
        an ACTION TITLE (a full-sentence assertion — the governing thought), then
        a thin rule. Long assertion titles wrap to two lines and drop the muted
        subtitle; short topic titles keep the subtitle as a descriptor."""
        mx = L["margin_x"]
        _kicker(slide, kicker, mx, L["margin_top"])
        long = len(str(title)) > 46
        tsize = 23 if long else 30
        _tb(slide, title, mx, L["margin_top"] + (0.1 if long else 0.14),
            TE.CANVAS["w"] - 2 * mx, (0.98 if long else 0.62),
            size=tsize, bold=True, color=col("title_text"), font=F_HEAD, spacing=1.04)
        if subtitle and not long:
            _tb(slide, subtitle, mx, L["margin_top"] + 0.86, TE.CANVAS["w"] - 2 * mx, 0.4,
                size=13.5, color=col("muted"), font=F_BODY)
        _rect(slide, mx, L["content_top"] - 0.16, TE.CANVAS["w"] - 2 * mx, 0.02,
              fill=col("hairline"))

    # ── per-layout renderers ──────────────────────────────────────────────
    def render_title(slide, s):
        _bg_gradient(slide, "bg3", angle=35)  # charcoal hero, subtle diagonal wash

        # Hero illustration: dedicated right-side region sized to ~40% of
        # slide width (within the 35-45% spec), full height minus the
        # footer band — placed after the gradient so it reads as blended
        # into the hero rather than pasted on top. Text column narrows to
        # make room; a glass card sits behind the title block so it stays
        # legible over any image content. Decorative blobs are skipped
        # entirely when a hero image is present — the image already
        # provides visual interest in that corner, and layering blobs on
        # top just adds clutter over the picture.
        hero_image = s.get("hero_image")
        has_hero = bool(hero_image and Path(hero_image).exists())
        if not has_hero:
            _decorative_blobs(slide, "accent", "on_dark", corner="br")
        text_w = 7.2 if has_hero else 11.5
        if has_hero:
            img_w = TE.CANVAS["w"] * 0.40
            _picture_cover(slide, hero_image, TE.CANVAS["w"] - img_w, 0, img_w, TE.CANVAS["h"] - 0.9)
            _card_glass(slide, 0.35, 2.3, text_w + 0.5, 3.9)

        # Yellow signature block, lower-left
        _rect(slide, 0, TE.CANVAS["h"] - 0.9, TE.CANVAS["w"], 0.9, fill=col("bg3"))
        _rect(slide, 0, 0, 0.16, TE.CANVAS["h"], fill=hero_accent())
        # brand mark
        _tb(slide, "EY", 0.65, 0.55, 2.0, 0.9, size=40, bold=True, color=hero_accent())
        _tb(slide, "Autonomous SDLC Studio", 1.7, 0.72, 6.0, 0.5, size=13,
            color=col("on_dark"), font=F_BODY)
        badge = s.get("badge", "")
        if badge:
            bw = 0.16 + 0.11 * len(badge)
            badge_x = TE.CANVAS["w"] - 0.55 - bw if not has_hero else TE.CANVAS["w"] * 0.6 - bw
            _icon = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, badge_x,
                           0.6, bw, 0.44, fill=col("accent"))
            _tb(slide, badge.upper(), badge_x, 0.66, bw, 0.34,
                size=11, bold=True, color=col("on_accent"), align=PP_ALIGN.CENTER)
        # Title block
        _rect(slide, 0.65, 2.55, 2.6, 0.07, fill=hero_accent())
        _tb(slide, s.get("title", project_name), 0.62, 2.75, text_w, 2.3, size=50 if not has_hero else 40,
            bold=True, color=col("on_dark"), font=F_HEAD, spacing=1.02)
        _tb(slide, s.get("subtitle", "AI-Autonomous Software Delivery"), 0.65, 5.15, text_w - 0.5, 0.9,
            size=19 if not has_hero else 16, color=hero_accent(), font=F_BODY, spacing=1.1)
        date_str = datetime.now().strftime("%B %Y")
        _tb(slide, f"CONFIDENTIAL   ·   {date_str}", 0.65, 6.55, 8.0, 0.4,
            size=11, color=RGBColor(*_rgb(theme["muted"])), font=F_BODY)
        _notes(slide, s.get("speaker_notes", ""))

    def render_section(slide, s):
        _bg_gradient(slide, "bg3", angle=135)

        hero_image = s.get("hero_image")
        has_hero = bool(hero_image and Path(hero_image).exists())
        if not has_hero:
            _decorative_blobs(slide, "accent", "on_dark", corner="tr")
        text_w = 7.2 if has_hero else 11.5
        if has_hero:
            img_w = TE.CANVAS["w"] * 0.40
            _picture_cover(slide, hero_image, TE.CANVAS["w"] - img_w, 0, img_w, TE.CANVAS["h"])
            _card_glass(slide, 0.4, 1.9, text_w + 0.4, 4.4)

        _rect(slide, 0, 0, 0.16, TE.CANVAS["h"], fill=hero_accent())
        num = s.get("number", "")
        if num:
            _tb(slide, str(num), 0.65, 2.2, 3.0, 1.6, size=110, bold=True,
                color=hero_accent(), font=F_HEAD)
        _rect(slide, 0.72, 4.15, 2.4, 0.06, fill=hero_accent())
        _tb(slide, s.get("title", ""), 0.68, 4.35, text_w, 1.4, size=42 if not has_hero else 34,
            bold=True, color=col("on_dark"), font=F_HEAD)
        if s.get("subtitle"):
            _tb(slide, s["subtitle"], 0.7, 5.7, text_w - 0.5, 0.7, size=16, color=hero_accent(),
                font=F_BODY)
        _notes(slide, s.get("speaker_notes", ""))

    def render_items(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Overview"))
        items = s.get("items", [])
        mx = L["margin_x"]
        top = L["content_top"]
        avail = L["content_bottom"] - top
        diagram_path = s.get("diagram_image")
        has_diagram = bool(diagram_path and Path(diagram_path).exists())
        # When an auto-generated diagram is attached, split the slide:
        # narrower bullet column on the left, diagram on the right — so the
        # diagram the user asked for ("generate diagrams wherever applicable")
        # is actually visible, not just generated and discarded.
        content_w = (TE.CANVAS["w"] - 2 * mx) * 0.52 if has_diagram else (TE.CANVAS["w"] - 2 * mx)
        n = max(1, min(len(items), 6))
        card_h = min((avail - (n - 1) * 0.16) / n, 0.95)
        y = top
        for it in items[:n]:
            _card(slide, mx, y, content_w, card_h)
            _rect(slide, mx, y, L["accent_bar_w"], card_h, fill=col("accent"))
            _icon_chip(slide, mx + 0.22, y + (card_h - 0.5) / 2, 0.5, it.get("icon", "check"))
            has_body = bool(it.get("body"))
            ty = y + (0.16 if has_body else (card_h - 0.4) / 2)
            _tb(slide, it.get("title", ""), mx + 0.95, ty, content_w - 1.2, 0.42,
                size=16 if not has_diagram else 13.5, bold=True, color=col("title_text"), font=F_HEAD)
            if has_body:
                _tb(slide, it["body"], mx + 0.95, ty + 0.42, content_w - 1.2, 0.4,
                    size=12.5, color=col("body_text"), font=F_BODY)
            y += card_h + 0.16
        if has_diagram:
            dx = mx + content_w + L["gutter"]
            dw = (TE.CANVAS["w"] - mx) - dx
            _card(slide, dx, top, dw, avail, fill="bg2")
            _picture(slide, diagram_path, dx + 0.15, top + 0.15, dw - 0.3, avail - 0.3)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_agenda(slide, s):
        _header(slide, s.get("title", "Agenda"), s.get("subtitle", ""), s.get("kicker", "Contents"))
        items = s.get("items", [])
        mx = L["margin_x"]
        top = L["content_top"]
        col_w = (TE.CANVAS["w"] - 2 * mx - L["gutter"]) / 2
        row_h = 0.92
        for i, it in enumerate(items[:6]):
            c = i % 2; r = i // 2
            x = mx + c * (col_w + L["gutter"]); y = top + r * (row_h + 0.14)
            _card(slide, x, y, col_w, row_h)
            _tb(slide, f"{i + 1:02d}", x + 0.22, y + 0.14, 0.9, 0.7, size=30, bold=True,
                color=RGBColor(*_rgb(theme["accent"])), font=F_HEAD)
            _rect(slide, x + 1.15, y + 0.22, 0.02, row_h - 0.44, fill=col("hairline"))
            _tb(slide, it.get("title", ""), x + 1.35, y + 0.16, col_w - 1.5, 0.42,
                size=15, bold=True, color=col("title_text"), font=F_HEAD)
            if it.get("body"):
                _tb(slide, it["body"], x + 1.35, y + 0.54, col_w - 1.5, 0.32,
                    size=11.5, color=col("muted"), font=F_BODY)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_kpi_cards(slide, s):
        _hero_wash(slide, s)
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "By the Numbers"))
        stats = s.get("stats", [])[:4]
        mx = L["margin_x"]
        n = max(1, len(stats))
        gap = L["gutter"]
        cw = (TE.CANVAS["w"] - 2 * mx - (n - 1) * gap) / n
        top = L["content_top"] + 0.35
        ch = 2.9
        for i, st in enumerate(stats):
            x = mx + i * (cw + gap)
            _card(slide, x, top, cw, ch)
            _rect(slide, x, top, cw, 0.09, fill=col("accent"))
            value_str = str(st.get("value", ""))
            # Adaptive size: a long value (e.g. "3.2x" vs "48h") must never
            # overflow the card and collide with the label below it.
            value_size = 58 if len(value_str) <= 4 else max(30, 58 - (len(value_str) - 4) * 5)
            _tb(slide, value_str, x, top + 0.55, cw, 1.15, size=value_size, bold=True,
                color=col("title_text" if is_dark else "accent2"), font=F_HEAD,
                align=PP_ALIGN.CENTER)
            _tb(slide, str(st.get("label", "")), x + 0.15, top + 1.75, cw - 0.3, 0.5,
                size=15, bold=True, color=col("title_text"), font=F_HEAD, align=PP_ALIGN.CENTER)
            if st.get("sub"):
                _tb(slide, st["sub"], x + 0.2, top + 2.25, cw - 0.4, 0.55, size=11.5,
                    color=col("muted"), font=F_BODY, align=PP_ALIGN.CENTER)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_stats_grid(slide, s):  # 2×2 variant
        _hero_wash(slide, s)
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "By the Numbers"))
        stats = s.get("stats", [])[:4]
        mx = L["margin_x"]; gap = L["gutter"]
        cw = (TE.CANVAS["w"] - 2 * mx - gap) / 2
        top = L["content_top"]
        rows = (len(stats) + 1) // 2
        ch = (L["content_bottom"] - top - (rows - 1) * gap) / max(rows, 1)
        for i, st in enumerate(stats):
            c = i % 2; r = i // 2
            x = mx + c * (cw + gap); y = top + r * (ch + gap)
            _card(slide, x, y, cw, ch)
            _rect(slide, x, y, L["accent_bar_w"], ch, fill=col("accent"))
            _tb(slide, str(st.get("value", "")), x + 0.35, y + 0.2, cw - 0.5, ch * 0.5,
                size=46, bold=True, color=col("title_text" if is_dark else "accent2"), font=F_HEAD)
            _tb(slide, str(st.get("label", "")), x + 0.38, y + ch * 0.58, cw - 0.6, 0.45,
                size=15, bold=True, color=col("title_text"), font=F_HEAD)
            if st.get("sub"):
                _tb(slide, st["sub"], x + 0.38, y + ch * 0.58 + 0.42, cw - 0.6, 0.4,
                    size=11.5, color=col("muted"), font=F_BODY)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_table(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Detail"))
        td = s.get("table", {})
        headers = td.get("headers", [])
        data_rows = td.get("rows", [])
        callout = s.get("callout")
        mx = L["margin_x"]
        top = L["content_top"]
        tbl_w = (TE.CANVAS["w"] - 2 * mx) if not callout else 8.4
        col_w = tbl_w / max(len(headers), 1)
        hdr_h = 0.52
        n_rows = min(len(data_rows), 11)
        # Adaptive row height fills the content band so the table never floats in
        # a sea of empty space (consulting decks use the whole canvas).
        avail = L["content_bottom"] - top - hdr_h
        row_h = min(max(avail / max(n_rows, 1), 0.46), 0.62)
        # header band
        _rect(slide, mx, top, tbl_w, hdr_h, fill=col("table_hdr"))
        for hi, h in enumerate(headers):
            _tb(slide, str(h), mx + hi * col_w + 0.14, top, col_w - 0.2, hdr_h,
                size=12, bold=True, color=col("table_hdr_text"), font=F_HEAD,
                anchor=MSO_ANCHOR.MIDDLE)
        for ri, row in enumerate(data_rows[:n_rows]):
            ry = top + hdr_h + ri * row_h
            if ri % 2 == 1:
                _rect(slide, mx, ry, tbl_w, row_h, fill=col("table_alt"))
            _rect(slide, mx, ry + row_h - 0.01, tbl_w, 0.01, fill=col("hairline"))
            for ci, cell in enumerate(row[:len(headers)]):
                cs = str(cell)
                cc = col("body_text")
                if cs.upper() in ("CRITICAL", "HIGH", "FAILED", "AT RISK"):
                    cc = col("danger")
                elif cs.upper() in ("PASSED", "COMPLIANT", "✓", "DONE", "COMPLETE"):
                    cc = col("success")
                _tb(slide, cs[:46], mx + ci * col_w + 0.14, ry, col_w - 0.22, row_h,
                    size=11.5, color=cc, font=F_BODY, bold=(ci == 0), anchor=MSO_ANCHOR.MIDDLE)
        if callout:
            _render_callout(slide, callout, tbl_w + mx + 0.35, top,
                            TE.CANVAS["w"] - mx - (tbl_w + mx + 0.35), L["content_bottom"] - top)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_chart(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Metrics"))
        cd = s.get("chart", {})
        data_items = cd.get("data", [])
        callout = s.get("callout")
        mx = L["margin_x"]
        chart_w = (8.2 if callout else TE.CANVAS["w"] - 2 * mx)
        label_w = 3.1
        track_x = mx + label_w + 0.15
        track_w = chart_w - label_w - 0.7
        n = min(len(data_items), 7)
        # Fill the content band and vertically centre the group: taller bars for
        # sparse charts, tighter for dense ones (consulting decks avoid dead space).
        band_top, band_bot = L["content_top"] + 0.15, L["content_bottom"] - 0.1
        band_h = band_bot - band_top
        row_h = min(max(band_h / max(n, 1), 0.5), 0.95)
        bar_h = min(row_h * 0.5, 0.42)
        group_h = row_h * n
        y = band_top + max((band_h - group_h) / 2, 0)
        for it in data_items[:n]:
            lbl = str(it.get("label", ""))[:34]
            val = float(it.get("value", 0)); mxv = float(it.get("max", 100)) or 100
            ratio = min(val / mxv, 1.0)
            cy = y + (row_h - bar_h) / 2
            _tb(slide, lbl, mx, y, label_w, row_h, size=13.5, color=col("body_text"),
                font=F_BODY, anchor=MSO_ANCHOR.MIDDLE)
            _rect(slide, track_x, cy, track_w, bar_h, fill=col("table_alt" if not is_dark else "bg3"))
            fillw = max(track_w * ratio, 0.08)
            _rect(slide, track_x, cy, fillw, bar_h, fill=col("accent"))
            vstr = f"{val:.0f}%" if mxv == 100 else f"{val:.0f}"
            _tb(slide, vstr, track_x + fillw + 0.12, y, 1.1, row_h, size=14, bold=True,
                color=col("accent2" if not is_dark else "accent"), font=F_HEAD,
                anchor=MSO_ANCHOR.MIDDLE)
            y += row_h
        if callout:
            _render_callout(slide, callout, chart_w + mx + 0.35, L["content_top"],
                            TE.CANVAS["w"] - mx - (chart_w + mx + 0.35), L["content_bottom"] - L["content_top"])
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def _render_callout(slide, callout, x, y, w, h):
        ctype = callout.get("type", "info")
        cmap = {"info": "info", "success": "success", "warning": "warning", "danger": "danger"}
        cc = col(cmap.get(ctype, "info"))
        _card(slide, x, y, w, h, fill="bg2")
        _rect(slide, x, y, w, 0.09, fill=cc)
        _tb(slide, callout.get("title", ""), x + 0.22, y + 0.24, w - 0.4, 0.5, size=14, bold=True,
            color=cc, font=F_HEAD)
        lines = [l.lstrip("•✓ ").strip() for l in str(callout.get("body", "")).split("\n") if l.strip()]
        _bullets(slide, lines, x + 0.22, y + 0.85, w - 0.44, h - 1.0, size=11.5,
                 marker="—", marker_color=cc)

    def render_two_col(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Comparison"))
        mx = L["margin_x"]; gap = L["gutter"]; top = L["content_top"]
        cw = (TE.CANVAS["w"] - 2 * mx - gap) / 2
        ch = L["content_bottom"] - top
        for idx, (hk, ck) in enumerate([("left_header", "left_content"), ("right_header", "right_content")]):
            x = mx + idx * (cw + gap)
            _card(slide, x, top, cw, ch)
            _rect(slide, x, top, cw, 0.55, fill=col("accent" if idx == 0 else "bg3"))
            _tb(slide, s.get(hk, ""), x + 0.25, top + 0.09, cw - 0.4, 0.4, size=15, bold=True,
                color=col("on_accent" if idx == 0 else "on_dark"), font=F_HEAD)
            lines = [l.strip() for l in str(s.get(ck, "")).split("\n") if l.strip()]
            _bullets(slide, lines, x + 0.28, top + 0.75, cw - 0.5, ch - 0.9, size=13,
                     marker="—", marker_color=col("accent"))
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_comparison(slide, s):
        """Before/after or option-A/option-B with a vs. divider."""
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Before / After"))
        mx = L["margin_x"]; top = L["content_top"]
        cw = (TE.CANVAS["w"] - 2 * mx - 1.0) / 2
        ch = L["content_bottom"] - top
        left = s.get("left", {}); right = s.get("right", {})
        for idx, (data, accent) in enumerate([(left, "danger"), (right, "success")]):
            x = mx + idx * (cw + 1.0)
            _card(slide, x, top, cw, ch)
            _rect(slide, x, top, L["accent_bar_w"], ch, fill=col(accent))
            _tb(slide, data.get("header", ""), x + 0.3, top + 0.22, cw - 0.5, 0.5, size=17,
                bold=True, color=col(accent), font=F_HEAD)
            lines = data.get("points", []) or [l for l in str(data.get("content", "")).split("\n") if l.strip()]
            _bullets(slide, lines, x + 0.32, top + 0.95, cw - 0.55, ch - 1.1, size=13,
                     marker=("✕" if idx == 0 else "✓"), marker_color=col(accent))
        # vs. chip in the middle
        _icon_chip(slide, mx + cw + 0.35, top + ch / 2 - 0.3, 0.6, "vs", fill="accent2",
                   glyph_color="on_dark")
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_tech_grid(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Technology"))
        items = s.get("tech_items", [])[:8]
        mx = L["margin_x"]; gap = L["gutter"]; top = L["content_top"]
        cols_n = 4 if len(items) > 4 else max(len(items), 1)
        rows_n = (len(items) + cols_n - 1) // cols_n
        cw = (TE.CANVAS["w"] - 2 * mx - (cols_n - 1) * gap) / cols_n
        ch = min((L["content_bottom"] - top - (rows_n - 1) * gap) / max(rows_n, 1), 2.3)
        for i, t in enumerate(items):
            c = i % cols_n; r = i // cols_n
            x = mx + c * (cw + gap); y = top + r * (ch + gap)
            _card(slide, x, y, cw, ch)
            _icon_chip(slide, x + 0.24, y + 0.24, 0.6, t.get("icon", "gear"))
            _tb(slide, str(t.get("layer", "")).upper(), x + 0.26, y + 1.0, cw - 0.4, 0.3,
                size=10, bold=True, color=col("accent3"), font=F_HEAD)
            _tb(slide, str(t.get("tech", "")), x + 0.26, y + 1.32, cw - 0.4, ch - 1.4, size=13.5,
                bold=True, color=col("title_text"), font=F_HEAD, spacing=1.05)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_process(slide, s):
        """Horizontal numbered process flow with arrows between steps."""
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "How It Works"))
        steps = s.get("steps", [])[:5]
        mx = L["margin_x"]; top = L["content_top"] + 0.4
        n = max(len(steps), 1)
        gap = 0.35
        cw = (TE.CANVAS["w"] - 2 * mx - (n - 1) * gap) / n
        ch = 3.3
        for i, st in enumerate(steps):
            x = mx + i * (cw + gap)
            _card(slide, x, top, cw, ch)
            _shape(slide, MSO_SHAPE.OVAL, x + cw / 2 - 0.4, top + 0.3, 0.8, 0.8, fill=col("accent"))
            _tb(slide, str(i + 1), x + cw / 2 - 0.4, top + 0.4, 0.8, 0.6, size=30, bold=True,
                color=col("on_accent"), font=F_HEAD, align=PP_ALIGN.CENTER)
            _tb(slide, st.get("title", ""), x + 0.18, top + 1.35, cw - 0.36, 0.7, size=14.5,
                bold=True, color=col("title_text"), font=F_HEAD, align=PP_ALIGN.CENTER)
            if st.get("body"):
                _tb(slide, st["body"], x + 0.18, top + 2.05, cw - 0.36, ch - 2.1, size=11.5,
                    color=col("body_text"), font=F_BODY, align=PP_ALIGN.CENTER, spacing=1.1)
            if i < n - 1:
                _tb(slide, "→", x + cw + 0.02, top + ch / 2 - 0.3, gap, 0.6, size=26, bold=True,
                    color=col("accent3"), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_architecture(slide, s):
        """Layered architecture: horizontal tiers stacked vertically."""
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Architecture"))
        layers = s.get("layers", [])[:5]
        mx = L["margin_x"]; top = L["content_top"]
        n = max(len(layers), 1)
        gap = 0.16
        lh = (L["content_bottom"] - top - (n - 1) * gap) / n
        for i, lyr in enumerate(layers):
            y = top + i * (lh + gap)
            _card(slide, mx, y, TE.CANVAS["w"] - 2 * mx, lh)
            _rect(slide, mx, y, 2.7, lh, fill=col("bg3"))
            _icon_chip(slide, mx + 0.22, y + lh / 2 - 0.27, 0.54, lyr.get("icon", "layers"))
            _tb(slide, lyr.get("name", ""), mx + 0.9, y + 0.12, 1.75, lh - 0.24, size=13.5,
                bold=True, color=col("on_dark"), font=F_HEAD, anchor=MSO_ANCHOR.MIDDLE)
            comps = lyr.get("components", [])
            if comps:
                cw2 = (TE.CANVAS["w"] - 2 * mx - 2.9 - 0.2) / max(len(comps), 1)
                for j, cptxt in enumerate(comps[:5]):
                    cx = mx + 2.85 + j * cw2
                    chip = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx + 0.06, y + 0.16,
                                  cw2 - 0.14, lh - 0.32, fill=col("bg2"), line_rgb=col("hairline"),
                                  line_w=0.75)
                    try:
                        chip.adjustments[0] = 0.12
                    except Exception:
                        pass
                    _tb(slide, str(cptxt), cx + 0.06, y + 0.16, cw2 - 0.14, lh - 0.32, size=11.5,
                        bold=True, color=col("title_text"), font=F_HEAD, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
            elif lyr.get("desc"):
                _tb(slide, lyr["desc"], mx + 2.95, y + 0.12, TE.CANVAS["w"] - 2 * mx - 3.1,
                    lh - 0.24, size=12.5, color=col("body_text"), font=F_BODY,
                    anchor=MSO_ANCHOR.MIDDLE)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_roadmap(slide, s):
        """Phase cards along a horizontal spine."""
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Roadmap"))
        phases = s.get("phases", s.get("timeline", []))[:5]
        mx = L["margin_x"]; top = L["content_top"] + 0.3
        n = max(len(phases), 1)
        gap = 0.3
        cw = (TE.CANVAS["w"] - 2 * mx - (n - 1) * gap) / n
        spine_y = top + 0.5
        _rect(slide, mx + cw / 2, spine_y, TE.CANVAS["w"] - 2 * mx - cw, 0.03, fill=col("hairline"))
        for i, ph in enumerate(phases):
            x = mx + i * (cw + gap)
            _shape(slide, MSO_SHAPE.OVAL, x + cw / 2 - 0.14, spine_y - 0.11, 0.28, 0.28, fill=col("accent"))
            _tb(slide, str(ph.get("date", ph.get("phase", ""))), x, top + 0.9, cw, 0.35, size=12,
                bold=True, color=col("accent3"), font=F_HEAD, align=PP_ALIGN.CENTER)
            _card(slide, x, top + 1.3, cw, 2.9)
            _rect(slide, x, top + 1.3, cw, 0.08, fill=col("accent"))
            _tb(slide, str(ph.get("title", "")), x + 0.16, top + 1.5, cw - 0.32, 0.7, size=14.5,
                bold=True, color=col("title_text"), font=F_HEAD, align=PP_ALIGN.CENTER)
            desc = ph.get("desc", ph.get("body", ""))
            if desc:
                _tb(slide, desc, x + 0.16, top + 2.2, cw - 0.32, 1.9, size=11.5,
                    color=col("body_text"), font=F_BODY, align=PP_ALIGN.CENTER, spacing=1.1)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_timeline(slide, s):
        events = s.get("timeline", [])
        if events and isinstance(events[0], dict) and events[0].get("desc") and len(events) <= 5:
            return render_roadmap(slide, {**s, "phases": events})
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", "Timeline"))
        n = min(len(events), 7)
        if n == 0:
            _footer(slide, s["_n"], s["_total"], s.get("title", "")); return
        mx = L["margin_x"]
        tl_y = 3.9
        start_x = mx + 0.4; end_x = TE.CANVAS["w"] - mx - 0.4
        step = (end_x - start_x) / max(n - 1, 1)
        _rect(slide, start_x, tl_y, end_x - start_x, 0.03, fill=col("hairline"))
        for ei, ev in enumerate(events[:n]):
            ex = start_x + ei * step
            up = ei % 2 == 0
            _shape(slide, MSO_SHAPE.OVAL, ex - 0.12, tl_y - 0.09, 0.24, 0.24, fill=col("accent"))
            ty = tl_y - 1.15 if up else tl_y + 0.35
            _tb(slide, str(ev.get("date", ""))[:12], ex - 1.0, ty, 2.0, 0.32, size=11.5, bold=True,
                color=col("accent2" if not is_dark else "accent"), font=F_HEAD, align=PP_ALIGN.CENTER)
            _tb(slide, str(ev.get("title", ""))[:18], ex - 1.0, ty + 0.32, 2.0, 0.35, size=13,
                bold=True, color=col("title_text"), font=F_HEAD, align=PP_ALIGN.CENTER)
            if ev.get("desc"):
                _tb(slide, str(ev["desc"])[:26], ex - 1.0, ty + 0.66, 2.0, 0.35, size=10.5,
                    color=col("muted"), font=F_BODY, align=PP_ALIGN.CENTER)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_quote(slide, s):
        _bg_gradient(slide, "bg3", angle=90)

        hero_image = s.get("hero_image")
        has_hero = bool(hero_image and Path(hero_image).exists())
        if not has_hero:
            _decorative_blobs(slide, "accent", "on_dark", corner="tr")
        quote_w = 6.6 if has_hero else 10.4
        if has_hero:
            img_w = TE.CANVAS["w"] * 0.38
            _picture_cover(slide, hero_image, TE.CANVAS["w"] - img_w, 0, img_w, TE.CANVAS["h"])
            _card_glass(slide, 0.4, 0.8, quote_w + 1.3, 4.7)

        _rect(slide, 0, 0, 0.16, TE.CANVAS["h"], fill=hero_accent())
        _tb(slide, "“", 0.7, 1.0, 3.0, 2.0, size=140, bold=True,
            color=hero_accent(), font=F_HEAD)
        _tb(slide, s.get("quote", s.get("title", "")), 1.9, 2.4, quote_w, 2.6, size=30 if not has_hero else 24,
            bold=True, color=col("on_dark"), font=F_HEAD, spacing=1.12)
        if s.get("attribution"):
            _rect(slide, 1.95, 5.3, 0.5, 0.05, fill=hero_accent())
            _tb(slide, s["attribution"], 2.6, 5.15, quote_w - 1.0, 0.5, size=15, color=hero_accent(),
                font=F_BODY)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_closing(slide, s):
        _bg_gradient(slide, "bg3", angle=45)

        # Hero image is confined to the top-right area only (above the icon-
        # items row, which needs the full slide width) — a smaller region
        # than title/section/quote's full-height treatment, but still gives
        # the closing slide visual presence instead of ending the deck on
        # a flat gradient.
        hero_image = s.get("hero_image")
        has_hero = bool(hero_image and Path(hero_image).exists())
        if not has_hero:
            _decorative_blobs(slide, "accent", "on_dark", corner="br")
        text_w = 7.2 if has_hero else 11.5
        if has_hero:
            img_w = TE.CANVAS["w"] * 0.35
            _picture_cover(slide, hero_image, TE.CANVAS["w"] - img_w, 0, img_w, 3.6)
            _card_glass(slide, 0.4, 0.35, text_w + 0.4, 3.0)

        _rect(slide, 0, 0, 0.16, TE.CANVAS["h"], fill=hero_accent())
        _tb(slide, "EY", 0.65, 0.5, 2.0, 0.8, size=32, bold=True, color=hero_accent())
        _rect(slide, 0.68, 1.7, 2.4, 0.06, fill=hero_accent())
        _tb(slide, s.get("title", "Thank You"), 0.65, 1.9, text_w, 1.2, size=44 if not has_hero else 36,
            bold=True, color=col("on_dark"), font=F_HEAD)
        if s.get("subtitle"):
            _tb(slide, s["subtitle"], 0.68, 3.05, text_w - 0.5, 0.5, size=15, color=hero_accent(),
                font=F_BODY)
        items = s.get("items", [])[:4]
        mx = 0.65; gap = 0.28
        if items:
            cw = (TE.CANVAS["w"] - 2 * mx - (len(items) - 1) * gap) / len(items)
            y = 3.9
            for i, it in enumerate(items):
                x = mx + i * (cw + gap)
                _icon_chip(slide, x, y, 0.5, it.get("icon", "check"))
                _tb(slide, it.get("title", ""), x, y + 0.65, cw, 0.6, size=13.5, bold=True,
                    color=col("on_dark"), font=F_HEAD)
                if it.get("body"):
                    _tb(slide, it["body"], x, y + 1.2, cw, 1.1, size=11, color=col("muted"),
                        font=F_BODY, spacing=1.1)
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    def render_hero(slide, s):
        """Full-bleed, image-led 'hero moment' slide — a single powerful
        statement/vision beat (Gamma-style), distinct from title/section/
        quote/closing where a hero image only ever decorates one side.
        Falls back to the standard hero gradient + decorative blobs when no
        hero_image is attached, so this layout never renders empty."""
        hero_image = s.get("hero_image")
        has_hero = bool(hero_image and Path(hero_image).exists())
        if has_hero:
            _picture_cover(slide, hero_image, 0, 0, TE.CANVAS["w"], TE.CANVAS["h"])
        else:
            _bg_gradient(slide, "bg3", angle=45)
            _decorative_blobs(slide, "accent", "on_dark", corner="tr")

        # Bottom-anchored glass card carries the statement over the image
        # (or gradient) so text stays legible either way.
        card_h = 2.9
        card_y = TE.CANVAS["h"] - card_h - 0.4
        _card_glass(slide, 0.6, card_y, TE.CANVAS["w"] - 1.2, card_h, tint="on_dark", alpha=55)
        _rect(slide, 0.9, card_y + 0.32, 0.5, 0.06, fill=hero_accent())
        _tb(slide, s.get("title", ""), 0.9, card_y + 0.5, TE.CANVAS["w"] - 2.2, 1.1,
            size=34, bold=True, color=col("on_dark"), font=F_HEAD, spacing=1.05)
        body = s.get("subtitle") or s.get("content") or ""
        if body:
            _tb(slide, body, 0.9, card_y + 1.6, TE.CANVAS["w"] - 2.2, 1.0,
                size=15, color=hero_accent(), font=F_BODY, spacing=1.15)
        _notes(slide, s.get("speaker_notes", ""))

    def render_content(slide, s):
        _header(slide, s.get("title", ""), s.get("subtitle", ""), s.get("kicker", ""))
        text = s.get("content", "") or s.get("subtitle", "")
        lines = [l.strip() for l in str(text).replace("•", "\n").split("\n") if l.strip()]
        mx = L["margin_x"]
        _card(slide, mx, L["content_top"], TE.CANVAS["w"] - 2 * mx, L["content_bottom"] - L["content_top"])
        _bullets(slide, lines, mx + 0.35, L["content_top"] + 0.3, TE.CANVAS["w"] - 2 * mx - 0.7,
                 L["content_bottom"] - L["content_top"] - 0.5, size=15, gap=8,
                 marker="—", marker_color=col("accent"))
        _footer(slide, s["_n"], s["_total"], s.get("title", ""))
        _notes(slide, s.get("speaker_notes", ""))

    RENDERERS = {
        "title": render_title, "section": render_section, "divider": render_section,
        "agenda": render_agenda, "items": render_items, "stats_grid": render_stats_grid,
        "kpi_cards": render_kpi_cards, "table": render_table, "chart": render_chart,
        "two_col": render_two_col, "two_column": render_two_col, "comparison": render_comparison,
        "tech_grid": render_tech_grid, "process": render_process,
        "architecture": render_architecture, "roadmap": render_roadmap,
        "timeline": render_timeline, "quote": render_quote, "closing": render_closing,
        "hero": render_hero, "content": render_content,
    }

    total = len(slides)
    for idx, s in enumerate(slides):
        s = dict(s)
        s["_n"] = idx + 1
        s["_total"] = total
        layout = s.get("layout", "content")
        if idx == 0 and layout not in ("title", "section"):
            layout = "title"
        slide = prs.slides.add_slide(blank)
        _bg(slide, "bg")
        renderer = RENDERERS.get(layout, render_content)
        try:
            renderer(slide, s)
            # Consulting-standard source citation (bottom-left, above footer).
            # Only on layouts that leave room at the bottom — full-height card
            # layouts (KPIs, columns, architecture) would otherwise overlap.
            src = s.get("source")
            if src and layout in ("chart", "table", "timeline", "roadmap", "process", "content"):
                _tb(slide, f"Source: {src}", L["margin_x"], L["footer_y"] - 0.34,
                    9.0, 0.3, size=9, italic=True, color=col("muted"), font=F_BODY)
        except Exception as exc:  # never let one bad slide kill the deck
            logger.warning("[PptxBuilder] layout '%s' failed on slide %d: %s", layout, idx + 1, exc)
            try:
                render_content(slide, s)
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info("[PptxBuilder] Built %d slides · theme=%s · %d bytes",
                total, theme.get("label"), buf.getbuffer().nbytes)
    return buf.getvalue()
